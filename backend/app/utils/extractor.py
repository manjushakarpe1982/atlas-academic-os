"""
extractor.py — Extract plain text from uploaded files.

Supports: PDF, DOCX, PPTX, TXT, MD.
Audio/video/image files return an empty string — they are handled
by separate transcription / vision pipelines (future phases).

All functions are synchronous and safe to call from a FastAPI
background task. They receive raw bytes (already read from storage)
and return a plain-text string.
"""

from __future__ import annotations
import io
from typing import Optional


# ── PDF ────────────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    try:
        import pypdf  # pip install pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text.strip())
        return "\n\n".join(pages)
    except Exception as e:
        raise RuntimeError(f"PDF extraction failed: {e}") from e


# ── DOCX ───────────────────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document  # pip install python-docx
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Also pull text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except Exception as e:
        raise RuntimeError(f"DOCX extraction failed: {e}") from e


# ── PPTX ───────────────────────────────────────────────────────────────────

def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # pip install python-pptx
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if len(parts) > 1:
                slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text)
    except Exception as e:
        raise RuntimeError(f"PPTX extraction failed: {e}") from e


# ── TXT / MD ───────────────────────────────────────────────────────────────

def _extract_text(data: bytes) -> str:
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# ── Public interface ────────────────────────────────────────────────────────

def extract_text(data: bytes, extension: str) -> str:
    """
    Extract plain text from file bytes.

    Args:
        data:      Raw file bytes.
        extension: Lowercase extension without dot (e.g. 'pdf', 'docx').

    Returns:
        Extracted text string, or empty string if the file type
        does not support text extraction (audio, video, image).

    Raises:
        RuntimeError: if a supported format fails to parse.
    """
    ext = extension.lower().lstrip(".")

    if ext == "pdf":
        return _extract_pdf(data)
    if ext in ("docx", "doc"):
        return _extract_docx(data)
    if ext in ("pptx", "ppt"):
        return _extract_pptx(data)
    if ext in ("txt", "md"):
        return _extract_text(data)

    # Audio, video, images — no text extraction in this phase
    return ""



def truncate_for_claude(text: str, max_chars: int = 12_000) -> str:
    """
    Truncate extracted text before sending to Claude.

    12,000 chars ≈ 3,000 tokens — enough for any syllabus or set of notes.
    Sending 80,000 chars (the old limit) wasted ~17,000 tokens per call.

    For syllabi the important content (grading, deadlines, topics) is almost
    always in the first half of the document. For notes/slides we send the
    first portion which covers the main topics.

    Override with ATLAS_MAX_PARSE_CHARS env var if you have very long docs
    and need higher accuracy.
    """
    import os
    max_chars = int(os.environ.get("ATLAS_MAX_PARSE_CHARS", max_chars))
    if len(text) <= max_chars:
        return text
    # Cut at a paragraph boundary where possible
    cut = text.rfind("\n\n", 0, max_chars)
    if cut > max_chars * 0.7:
        return text[:cut] + "\n\n[... document truncated for AI processing ...]"
    return text[:max_chars] + "\n\n[... document truncated for AI processing ...]"


def smart_truncate(text: str, category: str) -> str:
    """
    Category-aware truncation.
    Syllabi: take first 12k chars (grading info is always at the top).
    Slides:  spread across the whole doc to catch all topics.
    Notes:   first 10k is usually sufficient.
    Quiz:    typically short — 6k is plenty.
    """
    limits = {
        "syllabus":       12_000,
        "lecture_slides": 10_000,
        "notes":          10_000,
        "review_sheet":    8_000,
        "quiz":            6_000,
        "exam":            6_000,
        "graded_work":     6_000,
        "assignment":      8_000,
        "announcement":    4_000,
        "other":           8_000,
    }
    limit = limits.get(category, 8_000)
    return truncate_for_claude(text, limit)

